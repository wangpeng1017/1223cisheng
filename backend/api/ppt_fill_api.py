"""
@file ppt_fill_api.py
@desc PPT 表格填充 API - 接收半成品 PPT，填充表格数据，返回成品
"""
from fastapi import APIRouter, HTTPException, UploadFile, File
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
import shutil
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

router = APIRouter(prefix="/api/ppt/fill", tags=["PPT Fill"])

TEMP_DIR = Path("/tmp/ppt_fill")
TEMP_DIR.mkdir(parents=True, exist_ok=True)


class TableRow(BaseModel):
    values: List[str]


class TableData(BaseModel):
    slide_index: int  # 幻灯片索引（从1开始）
    table_index: int  # 表格索引（从0开始）
    header: Optional[List[str]] = None  # 表头
    body: List[TableRow]  # 表格数据


class FillTablesRequest(BaseModel):
    table_data: List[TableData]


def fill_table_in_ppt(ppt_path: str, table_data_list: List[TableData], output_path: str):
    """
    使用 python-pptx 填充 PPT 中的表格

    Args:
        ppt_path: 输入 PPT 文件路径
        table_data_list: 表格数据列表
        output_path: 输出 PPT 文件路径
    """
    if Presentation is None:
        raise HTTPException(status_code=500, detail="python-pptx 未安装")

    try:
        # 打开半成品 PPT
        prs = Presentation(ppt_path)

        # 遍历每个表格数据
        for table_data in table_data_list:
            slide_idx = table_data.slide_index - 1  # 转换为 0-based

            if slide_idx >= len(prs.slides):
                print(f"警告: 幻灯片索引 {table_data.slide_index} 超出范围")
                continue

            slide = prs.slides[slide_idx]
            table_idx = table_data.table_index

            # 查找表格
            tables = [shape for shape in slide.shapes if shape.has_table]
            if table_idx >= len(tables):
                print(f"警告: 幻灯片 {table_data.slide_index} 表格索引 {table_idx} 超出范围")
                continue

            table = tables[table_idx].table

            # 填充表头
            if table_data.header:
                if len(table.rows) < 1:
                    table.rows[0]
                for col_idx, header_text in enumerate(table_data.header):
                    if col_idx < len(table.columns):
                        cell = table.rows[0].cells[col_idx]
                        cell.text = str(header_text)

            # 填充表格数据
            for row_idx, row_data in enumerate(table_data.body):
                row_offset = 1 if table_data.header else 0

                # 如果需要更多行，添加新行
                while row_offset + row_idx >= len(table.rows):
                    table.rows._tbl.add_tr()

                # 填充单元格
                for col_idx, cell_value in enumerate(row_data.values):
                    if col_idx < len(table.columns):
                        cell = table.rows[row_offset + row_idx].cells[col_idx]
                        cell.text = str(cell_value)

        # 保存成品 PPT
        prs.save(output_path)
        return True

    except Exception as e:
        print(f"填充表格时出错: {str(e)}")
        raise HTTPException(status_code=500, detail=f"表格填充失败: {str(e)}")


@router.post("/tables")
async def fill_tables(
    file: UploadFile = File(...),
    table_data: str = None  # JSON 字符串
):
    """
    填充 PPT 中的表格

    Args:
        file: 半成品 PPT 文件（包含空表格）
        table_data: 表格数据（JSON 格式）
    """
    if not file.filename.endswith((".pptx", ".ppt")):
        raise HTTPException(status_code=400, detail="只支持 .pptx 或 .ppt 格式")

    # 保存上传的文件
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    input_path = TEMP_DIR / f"input_{timestamp}_{file.filename}"
    output_path = TEMP_DIR / f"output_{timestamp}_{file.filename}"

    try:
        # 保存上传的半成品 PPT
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        # 解析表格数据
        import json
        if table_data:
            data_dict = json.loads(table_data)
            table_data_list = [TableData(**td) for td in data_dict.get("table_data", [])]
        else:
            table_data_list = []

        # 填充表格
        fill_table_in_ppt(str(input_path), table_data_list, str(output_path))

        # 返回成品 PPT
        from fastapi.responses import FileResponse
        return FileResponse(
            path=str(output_path),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"filled_{file.filename}",
            background=None  # 不删除临时文件，由清理任务处理
        )

    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(status_code=500, detail=f"处理失败: {str(e)}")


@router.post("/tables/json")
async def fill_tables_json(request: FillTablesRequest):
    """
    填充 PPT 中的表格（JSON 格式请求）

    适用于已经知道 PPT 文件路径的场景
    """
    return {
        "message": "请使用 /api/ppt/fill/tables 端点上传文件",
        "example": "POST /api/ppt/fill/tables with file upload and table_data JSON"
    }
