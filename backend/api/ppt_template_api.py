"""
@file ppt_template_api.py
@desc PPT 模板 API 接口
"""
from fastapi import APIRouter, HTTPException, Depends, UploadFile, File
from sqlalchemy.orm import Session
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
from datetime import datetime
import os
import re
from pathlib import Path

from database import get_db
from models.ppt_template import PPTTemplate, PPTPlaceholder
from models.equipment import Equipment
from models.fixture import Fixture

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

router = APIRouter(prefix="/api/ppt/templates", tags=["PPT Templates"])

UPLOAD_DIR = Path("/root/npi-demo/backend/templates/uploaded")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


class TemplateResponse(BaseModel):
    id: int
    name: str
    entity_type: str
    file_path: Optional[str]
    slide_count: Optional[int]
    thumbnail_path: Optional[str]
    placeholders: List[Dict[str, Any]]
    created_at: datetime
    class Config:
        from_attributes = True


class PlaceholderConfig(BaseModel):
    placeholder: str
    field_key: str
    field_type: str
    default_value: Optional[str] = None


class BindTemplateRequest(BaseModel):
    template_id: int
    slide_index: int = 1


def extract_placeholders_from_text(text: str) -> List[str]:
    pattern = r"\[([^\]]+)\]"
    return re.findall(pattern, text)



def generate_thumbnail(file_path: str, output_dir: Path) -> str:
    """使用 LibreOffice 将 PPT 第一页转换为缩略图"""
    import subprocess
    import os
    
    output_dir.mkdir(parents=True, exist_ok=True)
    base_name = Path(file_path).stem
    thumbnail_path = output_dir / f"{base_name}_thumb.png"
    
    # 如果缩略图已存在，直接返回
    if thumbnail_path.exists():
        return str(thumbnail_path)
    
    # 使用 LibreOffice 转换为 PDF
    temp_dir = Path("/tmp/ppt_convert")
    temp_dir.mkdir(exist_ok=True)
    
    try:
        # 转换为 PDF
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(temp_dir),
            str(file_path)
        ], check=True, timeout=30, capture_output=True)
        
        pdf_file = temp_dir / f"{base_name}.pdf"
        
        # 使用 pdftoppm 或 convert 将 PDF 第一页转为图片
        img_path = temp_dir / f"{base_name}-1.png"
        
        # 尝试使用 pdftoppm (poppler-utils)
        try:
            subprocess.run([
                "pdftoppm", "-png", "-f", "1", "-singlefile",
                str(pdf_file), str(temp_dir / base_name)
            ], check=True, timeout=10, capture_output=True)
            if img_path.exists():
                img_path.rename(thumbnail_path)
                return str(thumbnail_path)
        except (FileNotFoundError, subprocess.CalledProcessError):
            pass
        
        # 尝试使用 convert (ImageMagick)
        try:
            subprocess.run([
                "convert", f"{pdf_file}[0]", str(thumbnail_path)
            ], check=True, timeout=10, capture_output=True)
            if thumbnail_path.exists():
                return str(thumbnail_path)
        except (FileNotFoundError, subprocess.CalledProcessError):
            pass
        
    except Exception as e:
        print(f"Thumbnail generation failed: {e}")
    finally:
        # 清理临时文件
        for f in temp_dir.glob("*"):
            try:
                f.unlink()
            except:
                pass
    
    return ""


def parse_ppt_template(file_path: str) -> Dict[str, Any]:
    if Presentation is None:
        return {"slide_count": 0, "slides": [], "detected_placeholders": []}
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        slides_info = []
        all_placeholders = set()
        for idx, slide in enumerate(prs.slides, start=1):
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
            for text in text_content:
                placeholders = extract_placeholders_from_text(text)
                all_placeholders.update(placeholders)
            slides_info.append({
                "index": idx,
                "text_preview": " ".join(text_content)[:200] if text_content else "",
                "placeholders": list(extract_placeholders_from_text(" ".join(text_content)))
            })
        return {
            "slide_count": slide_count,
            "slides": slides_info,
            "detected_placeholders": list(all_placeholders)
        }
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPT解析失败: {str(e)}")


@router.post("/upload")
async def upload_template(
    file: UploadFile = File(...),
    name: str = None,
    entity_type: str = "equipment",
    db: Session = Depends(get_db)
):
    if entity_type not in ["equipment", "fixture"]:
        raise HTTPException(status_code=400, detail="entity_type必须是equipment或fixture")
    if not file.filename.endswith((".pptx", ".ppt")):
        raise HTTPException(status_code=400, detail="只支持.pptx或.ppt格式")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{entity_type}_{timestamp}_{file.filename}"
    file_path = UPLOAD_DIR / filename
    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)
    ppt_info = parse_ppt_template(str(file_path))
    template_name = name or file.filename.replace(".pptx", "").replace(".ppt", "")
    
    # 生成缩略图
    thumbnail_dir = UPLOAD_DIR / "thumbnails"
    thumbnail_path = generate_thumbnail(str(file_path), thumbnail_dir)
    
    db_template = PPTTemplate(
        name=template_name,
        entity_type=entity_type,
        file_path=str(file_path),
        slide_count=ppt_info["slide_count"],
        placeholders=[{"key": p, "default": f"[{p}]"} for p in ppt_info["detected_placeholders"]],
        thumbnail_path=thumbnail_path
    )
    db.add(db_template)
    db.commit()
    db.refresh(db_template)
    for placeholder in ppt_info["detected_placeholders"]:
        db_placeholder = PPTPlaceholder(
            template_id=db_template.id,
            placeholder=f"[{placeholder}]",
            field_key=placeholder.lower().replace(" ", "_"),
            field_type="text"
        )
        db.add(db_placeholder)
    db.commit()
    return db_template


@router.get("")
def list_templates(
    entity_type: Optional[str] = None,
    skip: int = 0,
    limit: int = 100,
    db: Session = Depends(get_db)
):
    query = db.query(PPTTemplate)
    if entity_type:
        query = query.filter(PPTTemplate.entity_type == entity_type)
    return query.offset(skip).limit(limit).all()


@router.get("/{template_id}")
def get_template(template_id: int, db: Session = Depends(get_db)):
    template = db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    return template


@router.get("/{template_id}/slides")
def get_template_slides(template_id: int, db: Session = Depends(get_db)):
    template = db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    if not template.file_path or not os.path.exists(template.file_path):
        return {"slide_count": 0, "slides": []}
    ppt_info = parse_ppt_template(template.file_path)
    return ppt_info


@router.delete("/{template_id}")
def delete_template(template_id: int, db: Session = Depends(get_db)):
    template = db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    if template.file_path and os.path.exists(template.file_path):
        os.remove(template.file_path)
    db.delete(template)
    db.commit()
    return {"message": "模板已删除"}


@router.put("/{template_id}/placeholders")
def update_placeholders(
    template_id: int,
    placeholders: List[PlaceholderConfig],
    db: Session = Depends(get_db)
):
    template = db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    db.query(PPTPlaceholder).filter(PPTPlaceholder.template_id == template_id).delete()
    for p in placeholders:
        db_placeholder = PPTPlaceholder(
            template_id=template_id,
            placeholder=p.placeholder,
            field_key=p.field_key,
            field_type=p.field_type,
            default_value=p.default_value
        )
        db.add(db_placeholder)
    db.commit()
    return {"message": f"已更新{len(placeholders)}个占位符映射"}


@router.get("/{template_id}/placeholders")
def get_placeholders(template_id: int, db: Session = Depends(get_db)):
    placeholders = db.query(PPTPlaceholder).filter(PPTPlaceholder.template_id == template_id).all()
    return [{"placeholder": p.placeholder, "field_key": p.field_key,
             "field_type": p.field_type, "default_value": p.default_value} for p in placeholders]


@router.put("/equipment/{equipment_id}/template")
def bind_equipment_template(
    equipment_id: int,
    request: BindTemplateRequest,
    db: Session = Depends(get_db)
):
    equipment = db.query(Equipment).filter(Equipment.id == equipment_id).first()
    if not equipment:
        raise HTTPException(status_code=404, detail="设备不存在")
    template = db.query(PPTTemplate).filter(PPTTemplate.id == request.template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    if template.entity_type != "equipment":
        raise HTTPException(status_code=400, detail="模板类型不匹配")
    equipment.ppt_template_id = request.template_id
    equipment.ppt_slide_index = request.slide_index
    db.commit()
    return {"message": "模板绑定成功"}


@router.put("/fixture/{fixture_id}/template")
def bind_fixture_template(
    fixture_id: int,
    request: BindTemplateRequest,
    db: Session = Depends(get_db)
):
    fixture = db.query(Fixture).filter(Fixture.id == fixture_id).first()
    if not fixture:
        raise HTTPException(status_code=404, detail="夹具不存在")
    template = db.query(PPTTemplate).filter(PPTTemplate.id == request.template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    if template.entity_type != "fixture":
        raise HTTPException(status_code=400, detail="模板类型不匹配")
    fixture.ppt_template_id = request.template_id
    fixture.ppt_slide_index = request.slide_index
    db.commit()
    return {"message": "模板绑定成功"}


@router.get("/{template_id}/thumbnail")
def get_thumbnail(template_id: int, db: Session = Depends(get_db)):
    """获取模板缩略图"""
    from fastapi.responses import FileResponse
    
    template = db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    
    if not template.thumbnail_path or not os.path.exists(template.thumbnail_path):
        raise HTTPException(status_code=404, detail="缩略图不存在")
    
    return FileResponse(template.thumbnail_path, media_type="image/png")

@router.get("/{template_id}/download")
def download_template(template_id: int, db: Session = Depends(get_db)):
    """下载模板文件"""
    from fastapi.responses import FileResponse
    
    template = db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    
    if not template.file_path or not os.path.exists(template.file_path):
        raise HTTPException(status_code=404, detail="文件不存在")
    
    return FileResponse(
        template.file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=template.name + ".pptx"
    )

