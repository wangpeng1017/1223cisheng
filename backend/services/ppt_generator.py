"""
@file ppt_generator.py
@desc MTD PPT 自动生成服务 - 支持设备/夹具自定义模板
@see PRD: docs/mtd/PRD.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from datetime import datetime
from typing import List, Dict, Any, Optional
from copy import deepcopy
from sqlalchemy.orm import Session

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), '..', 'templates', 'mtd_template.pptx')
UPLOAD_DIR = os.path.join(os.path.dirname(__file__), '..', 'uploads', 'ppt_templates')


class MTDPPTGenerator:
    """MTD PPT 生成器 - 支持自定义设备/夹具模板"""

    def __init__(self, template_path: str = None, db: Session = None):
        self.template_path = template_path or TEMPLATE_PATH
        self.db = db
        self.prs = None

    def generate(self, project: Dict, equipment_list: List[Dict],
                 fixture_list: List[Dict], fai_items: List[Dict],
                 output_path: str) -> str:
        """
        生成 MTD PPT

        PPT 结构：
        - 第1页: 封面
        - 第2页: 目录
        - 第3页: 修订历史
        - 第4-N页: 设备详情（使用自定义模板或默认模板）
        - 夹具清单页
        - 测量详情总表
        - 测试项详情页
        """
        self.prs = Presentation(self.template_path)

        # 1. 更新封面页
        self._update_cover_slide(project)

        # 2. 更新目录页
        self._update_content_slide(project)

        # 3. 更新修订历史
        self._update_revision_slide(project)

        # 4. 处理设备详情页 - 支持自定义模板
        self._generate_equipment_slides(equipment_list, project)

        # 5. 处理夹具清单页 - 支持自定义模板
        self._generate_fixture_slides(fixture_list, project)

        # 6. 更新测量详情总表
        self._update_summary_table(fai_items)

        # 7. 更新测试项详情页
        self._update_detail_slides(fai_items, project)

        self.prs.save(output_path)
        return output_path

    def _update_cover_slide(self, project: Dict):
        """更新封面页（第1页）"""
        slide = self.prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Project Name" in text or "J510" in text or "Part Number" in text:
                    new_text = f"Project Name:{project.get('project_name', '')}\n"
                    new_text += f"Part Number & Rev：{project.get('part_number', '')}\n"
                    new_text += f"Vendor：{project.get('vendor', '')}\n"
                    new_text += f"Date：{datetime.now().strftime('%Y/%m/%d')}"
                    self._set_text_keep_format(shape, new_text)

    def _update_content_slide(self, project: Dict):
        """更新目录页（第2页）"""
        slide = self.prs.slides[1]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "MTD" in text and "Content" in text:
                    new_text = f"MTD | {project.get('project_name', '')}-Content-{project.get('part_number', '')}"
                    self._set_text_keep_format(shape, new_text)

    def _update_revision_slide(self, project: Dict):
        """更新修订历史页（第3页）"""
        slide = self.prs.slides[2]

        # 更新标题
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "MTD" in text and "revision" in text.lower():
                    new_text = f"MTD | {project.get('project_name', '')}- revision history List  {project.get('part_number', '')}"
                    self._set_text_keep_format(shape, new_text)

            # 更新表格
            if shape.has_table:
                table = shape.table
                if len(table.rows) > 1:
                    row = table.rows[1]
                    part_num = project.get('part_number', '')
                    row.cells[0].text = part_num.split('-')[0] if '-' in part_num else part_num
                    row.cells[1].text = project.get('revision', '01')
                    row.cells[2].text = datetime.now().strftime('%Y/%m/%d')
                    row.cells[3].text = "Initial Release"

    def _generate_equipment_slides(self, equipment_list: List[Dict], project: Dict):
        """生成设备详情页 - 支持自定义模板"""
        # 找到设备详情页的起始位置（第4页，索引3）
        # 先删除原有的设备页（第4-9页）
        while len(self.prs.slides) > 3:
            # 检查第4页是否还是设备页（有表格且包含设备相关字段）
            if len(self.prs.slides) > 3:
                slide = self.prs.slides[3]
                is_equipment_slide = False
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if 'name' in cell.text.lower() or 'manufacturer' in cell.text.lower():
                                    is_equipment_slide = True
                            if is_equipment_slide:
                                break
                    if is_equipment_slide:
                        break
                if is_equipment_slide:
                    # 删除这个占位设备页 - 使用正确的方式删除幻灯片
                    self._delete_slide(3)
                else:
                    break
            else:
                break

        # 为每个设备插入自定义模板页或默认页
        insert_pos = 3  # 第4页位置
        for equipment in equipment_list:
            template_id = equipment.get('ppt_template_id')
            slide_index = equipment.get('ppt_slide_index', 1)

            if template_id:
                # 使用自定义模板
                self._insert_custom_equipment_slide(equipment, template_id, slide_index, insert_pos, project)
            else:
                # 使用默认模板页
                self._insert_default_equipment_slide(equipment, insert_pos)
            insert_pos += 1

    def _insert_custom_equipment_slide(self, equipment: Dict, template_id: int, slide_index: int, insert_pos: int, project: Dict):
        """插入自定义设备模板页"""
        from models.ppt_template import PPTTemplate

        if not self.db:
            # 无数据库连接时使用默认方式
            self._insert_default_equipment_slide(equipment, insert_pos)
            return

        template = self.db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
        if not template or not template.file_path:
            self._insert_default_equipment_slide(equipment, insert_pos)
            return

        template_path = os.path.join(UPLOAD_DIR, template.file_path) if not os.path.isabs(template.file_path) else template.file_path

        if not os.path.exists(template_path):
            self._insert_default_equipment_slide(equipment, insert_pos)
            return

        try:
            # 加载模板 PPT
            template_prs = Presentation(template_path)

            # 获取指定页（slide_index 是 1-based）
            if 0 < slide_index <= len(template_prs.slides):
                source_slide = template_prs.slides[slide_index - 1]

                # 复制该页到主 PPT
                self._copy_slide(source_slide, self.prs, insert_pos)

                # 替换占位符
                new_slide = self.prs.slides[insert_pos]
                self._replace_placeholders(new_slide, equipment, project)
            else:
                self._insert_default_equipment_slide(equipment, insert_pos)
        except Exception as e:
            print(f"插入自定义模板失败: {e}")
            self._insert_default_equipment_slide(equipment, insert_pos)

    def _insert_default_equipment_slide(self, equipment: Dict, insert_pos: int):
        """插入默认设备页（从基础模板复制）"""
        # 如果还有默认设备页可用，复制一页
        if len(self.prs.slides) > 3:
            # 复制第4页作为模板
            template_slide = self.prs.slides[3]
            self._copy_slide(template_slide, self.prs, insert_pos)

            # 更新设备信息
            new_slide = self.prs.slides[insert_pos]
            for shape in new_slide.shapes:
                if shape.has_table:
                    table = shape.table
                    self._update_equipment_table(table, equipment)

    def _generate_fixture_slides(self, fixture_list: List[Dict], project: Dict):
        """生成夹具清单页"""
        # 找到夹具页位置
        fixture_slide_idx = self._find_fixture_slide_position()
        if fixture_slide_idx is None:
            return

        # 如果有夹具使用自定义模板，为每个夹具单独生成页面
        # 否则使用默认的清单页
        has_custom_template = any(f.get('ppt_template_id') for f in fixture_list)

        if has_custom_template:
            # 删除默认夹具页 - 使用正确的方式删除幻灯片
            self._delete_slide(fixture_slide_idx)

            # 为每个夹具插入自定义模板页
            insert_pos = fixture_slide_idx
            for fixture in fixture_list:
                template_id = fixture.get('ppt_template_id')
                slide_index = fixture.get('ppt_slide_index', 1)

                if template_id:
                    self._insert_custom_fixture_slide(fixture, template_id, slide_index, insert_pos, project)
                else:
                    self._insert_default_fixture_slide(fixture, insert_pos)
                insert_pos += 1
        else:
            # 使用默认夹具清单页
            slide = self.prs.slides[fixture_slide_idx]
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for idx, fixture in enumerate(fixture_list[:min(len(fixture_list), len(table.rows)-1)]):
                        row_idx = idx + 1
                        if row_idx < len(table.rows):
                            row = table.rows[row_idx]
                            if len(row.cells) >= 5:
                                row.cells[0].text = fixture.get('fixture_no', '')
                                row.cells[1].text = fixture.get('size', '')
                                row.cells[2].text = fixture.get('material', '')
                                row.cells[4].text = fixture.get('remark', '/')

    def _insert_custom_fixture_slide(self, fixture: Dict, template_id: int, slide_index: int, insert_pos: int, project: Dict):
        """插入自定义夹具模板页"""
        from models.ppt_template import PPTTemplate

        if not self.db:
            self._insert_default_fixture_slide(fixture, insert_pos)
            return

        template = self.db.query(PPTTemplate).filter(PPTTemplate.id == template_id).first()
        if not template or not template.file_path:
            self._insert_default_fixture_slide(fixture, insert_pos)
            return

        template_path = os.path.join(UPLOAD_DIR, template.file_path) if not os.path.isabs(template.file_path) else template.file_path

        if not os.path.exists(template_path):
            self._insert_default_fixture_slide(fixture, insert_pos)
            return

        try:
            template_prs = Presentation(template_path)

            if 0 < slide_index <= len(template_prs.slides):
                source_slide = template_prs.slides[slide_index - 1]
                self._copy_slide(source_slide, self.prs, insert_pos)

                new_slide = self.prs.slides[insert_pos]
                self._replace_placeholders(new_slide, fixture, project)
            else:
                self._insert_default_fixture_slide(fixture, insert_pos)
        except Exception as e:
            print(f"插入夹具自定义模板失败: {e}")
            self._insert_default_fixture_slide(fixture, insert_pos)

    def _insert_default_fixture_slide(self, fixture: Dict, insert_pos: int):
        """插入默认夹具页"""
        # 简单处理：插入一个基本信息页
        slide_layout = self.prs.slide_layouts[0]  # 使用空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        # 移动到正确位置
        self._move_slide_to_position(len(self.prs.slides) - 1, insert_pos)

    def _find_fixture_slide_position(self) -> Optional[int]:
        """找到夹具页的位置"""
        for idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if 'fixture' in shape.text.lower() or '夹具' in shape.text:
                        return idx
        return None

    def _delete_slide(self, slide_index: int):
        """删除指定索引的幻灯片"""
        if 0 <= slide_index < len(self.prs.slides):
            slide = self.prs.slides[slide_index]
            # 获取幻灯片的 rId 元素（不是整数 ID）
            rId = self.prs.slides._sldIdLst[slide_index]
            self.prs.part.drop_rel(rId.rId)
            self.prs.slides._sldIdLst.remove(rId)

    def _move_slide_to_position(self, from_index: int, to_index: int):
        """将幻灯片从一个位置移动到另一个位置"""
        if from_index == to_index:
            return
        if 0 <= from_index < len(self.prs.slides) and 0 <= to_index < len(self.prs.slides):
            # 获取要移动的幻灯片元素
            slide_elem = self.prs.slides._sldIdLst[from_index]
            # 从原位置移除
            self.prs.slides._sldIdLst.remove(slide_elem)
            # 插入到新位置
            self.prs.slides._sldIdLst.insert(to_index, slide_elem)

    def _copy_slide(self, source_slide, target_prs, insert_pos: int):
        """复制幻灯片到目标演示文稿"""
        # 创建新幻灯片
        slide_layout = target_prs.slide_layouts[0]  # 使用空白布局
        new_slide = target_prs.slides.add_slide(slide_layout)

        # 复制所有形状
        for shape in source_slide.shapes:
            self._copy_shape(shape, new_slide)

        # 移动到正确位置
        if insert_pos < len(target_prs.slides) - 1:
            # 获取新添加的幻灯片元素（在最后位置）
            slide_elem = target_prs.slides._sldIdLst[-1]
            target_prs.slides._sldIdLst.remove(slide_elem)
            target_prs.slides._sldIdLst.insert(insert_pos, slide_elem)

    def _copy_shape(self, shape, target_slide):
        """复制形状到目标幻灯片"""
        el = shape.element
        new_el = deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    def _replace_placeholders(self, slide, data: Dict, project: Dict):
        """替换幻灯片中的占位符"""
        # 合并数据和项目信息
        all_data = {**data, **project}

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                new_text = text

                # 替换 [字段名] 格式的占位符
                for key, value in all_data.items():
                    if value is not None:
                        placeholder = f'[{key}]'
                        if placeholder in new_text:
                            new_text = new_text.replace(placeholder, str(value))

                # 替换 specs.range 等嵌套字段
                if isinstance(data.get('specs'), dict):
                    for key, value in data['specs'].items():
                        placeholder = f'specs.{key}]'
                        if placeholder in new_text:
                            new_text = new_text.replace(placeholder, str(value))
                        # 也尝试直接字段名
                        placeholder2 = f'[{key}]'
                        if placeholder2 in new_text:
                            new_text = new_text.replace(placeholder2, str(value))

                if new_text != text:
                    self._set_text_keep_format(shape, new_text)

    def _update_equipment_table(self, table, equipment: Dict):
        """更新单个设备表格"""
        specs = equipment.get('specs', {}) or {}

        for row in table.rows:
            if len(row.cells) >= 2:
                label = row.cells[0].text.strip().lower()
                if 'name' in label:
                    row.cells[1].text = equipment.get('name', '')
                elif 'manufacturer' in label:
                    row.cells[1].text = equipment.get('manufacturer', '')
                elif 'model' in label:
                    row.cells[1].text = equipment.get('model', '')
                elif 'range' in label:
                    row.cells[1].text = specs.get('range', '')
                elif 'resolution' in label:
                    row.cells[1].text = specs.get('resolution', '')
                elif 'accuracy' in label:
                    row.cells[1].text = specs.get('accuracy', '')

    def _update_summary_table(self, fai_items: List[Dict]):
        """更新测量详情总表（第11页）"""
        if len(self.prs.slides) < 11:
            return

        slide = self.prs.slides[10]  # 第11页，索引10

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # 检查是否是主表格（10列）
                if len(table.columns) >= 10:
                    # 从第2行开始填充数据
                    for idx, item in enumerate(fai_items[:min(len(fai_items), len(table.rows)-1)]):
                        row_idx = idx + 1
                        if row_idx < len(table.rows):
                            row = table.rows[row_idx]
                            row.cells[0].text = str(item.get('fai_num') or '')
                            row.cells[1].text = item.get('spc') or ''
                            row.cells[2].text = item.get('specification') or ''
                            row.cells[3].text = item.get('description') or ''
                            row.cells[4].text = item.get('cpk_method') or ''
                            row.cells[5].text = item.get('cpk_fixture') or 'No'
                            row.cells[6].text = item.get('inprocess_method') or ''
                            row.cells[7].text = item.get('inprocess_fixture') or 'No'
                            row.cells[8].text = item.get('location') or ''
                            row.cells[9].text = item.get('cross_check_by') or ''

    def _update_detail_slides(self, fai_items: List[Dict], project: Dict):
        """更新测试项详情页（第12页起）"""
        # 详情页从索引11开始（第12页）
        detail_start_idx = 11

        for idx, item in enumerate(fai_items):
            slide_idx = detail_start_idx + idx
            if slide_idx >= len(self.prs.slides):
                break

            slide = self.prs.slides[slide_idx]

            # 更新页面标题
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "MTD" in text and "Metrology Details" in text:
                        new_title = f"MTD | {project.get('project_name', '')} Metrology Details-{project.get('part_number', '')}"
                        self._set_text_keep_format(shape, new_title)

                # 更新详情表格
                if shape.has_table:
                    table = shape.table
                    self._update_detail_table(table, item)

    def _update_detail_table(self, table, item: Dict):
        """更新单个测试项详情表格"""
        # 详情表格结构：3行2列
        # 第1行：标题（测试项名称和规格）
        # 第2行：CPK | In-process
        # 第3行：测量方法详情

        if len(table.rows) >= 1:
            # 更新标题行
            header_text = f"{item.get('description') or ''}  FAI {item.get('fai_num') or ''} /SPC {item.get('spc') or ''} :{item.get('specification') or ''}"
            if len(table.rows[0].cells) >= 1:
                table.rows[0].cells[0].text = header_text

        if len(table.rows) >= 3:
            # 更新测量方法详情
            cpk_text = f"1.  Measurement Method ：{item.get('cpk_method') or ''}\n"
            cpk_text += f"2.  Fixture ：{item.get('cpk_fixture', '/')}\n"
            cpk_text += f"3.  Measurement steps: {item.get('measurement_steps', 'See image')}"

            inprocess_text = f"1.  Measurement Method ：{item.get('inprocess_method') or ''}\n"
            inprocess_text += f"2.  Fixture ：{item.get('inprocess_fixture', '/')}\n"
            inprocess_text += f"3.  Measurement steps: {item.get('measurement_steps', 'See image')}"

            if len(table.rows[2].cells) >= 2:
                table.rows[2].cells[0].text = cpk_text
                table.rows[2].cells[1].text = inprocess_text

    def _set_text_keep_format(self, shape, new_text: str):
        """设置文本，尽量保持原有格式"""
        if shape.has_text_frame:
            # 保存第一个段落的格式
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    font = para.runs[0].font
                    font_name = font.name
                    font_size = font.size
                    font_bold = font.bold

                    # 清除并设置新文本
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = new_text

                    # 恢复格式
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if font_bold is not None:
                        run.font.bold = font_bold
                else:
                    shape.text_frame.text = new_text
            else:
                shape.text_frame.text = new_text


def generate_mtd_ppt(project: Dict, equipment_list: List[Dict],
                     fixture_list: List[Dict], fai_items: List[Dict],
                     output_path: str, db: Session = None) -> str:
    """便捷函数：生成 MTD PPT"""
    generator = MTDPPTGenerator(db=db)
    return generator.generate(project, equipment_list, fixture_list, fai_items, output_path)


def extract_images_from_ppt(ppt_path: str, output_dir: str) -> List[str]:
    """从 PPT 中提取所有图片"""
    prs = Presentation(ppt_path)
    os.makedirs(output_dir, exist_ok=True)

    image_paths = []
    img_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_count += 1
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext

                filename = f"slide{slide_idx+1}_img{img_count}.{image_ext}"
                filepath = os.path.join(output_dir, filename)

                with open(filepath, 'wb') as f:
                    f.write(image_bytes)

                image_paths.append(filepath)

    return image_paths
